import os
import logging
from flask import Flask, request, render_template, send_file
import pdfplumber
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import io
from transformers import AutoTokenizer, T5ForConditionalGeneration
import language_tool_python  # Import language tool

# Initialize Flask app
app = Flask(__name__)

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# Load Grammarly's CoEdit model and tokenizer
tokenizer = AutoTokenizer.from_pretrained("grammarly/coedit-large")
model = T5ForConditionalGeneration.from_pretrained("grammarly/coedit-large")

# Initialize LanguageTool for English
tool = language_tool_python.LanguageTool('en-US')

def correct_spelling_with_tool(text):
    """Correct spelling using LanguageTool."""
    try:
        logging.debug(f"Original text for spelling correction: {text}")
        
        # Check for spelling and grammar mistakes
        matches = tool.check(text)
        corrected_text = language_tool_python.utils.correct(text, matches)
        
        logging.debug(f"Corrected text after spelling correction: {corrected_text}")
        return corrected_text
    except Exception as e:
        logging.error(f"Error during spelling correction: {e}", exc_info=True)
        return "Error during spelling correction."

def correct_grammar_with_transformers(text):
    """Correct grammar using Grammarly's CoEdit model."""
    try:
        logging.debug(f"Original text for grammar correction: {text}")
        
        # Format input for the model
        input_text = f"Fix grammatical errors in this sentence: {text}"
        input_ids = tokenizer(input_text, return_tensors="pt").input_ids
        
        # Generate corrected text
        outputs = model.generate(input_ids, max_length=256)
        corrected_text = tokenizer.decode(outputs[0], skip_special_tokens=True)
        
        logging.debug(f"Corrected text: {corrected_text}")
        return corrected_text.strip()
    except Exception as e:
        logging.error(f"Error during grammar correction: {e}", exc_info=True)
        return "Error during grammar correction."

def extract_text_from_pdf(file):
    """Extract text from PDF file."""
    text = ''
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ''  # Handle None case
    logging.debug(f"Extracted PDF text: {text}")
    return text

def extract_text_from_ppt(file):
    """Extract text from PPTX file."""
    presentation = Presentation(file)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + ' '  # Add space between texts
    logging.debug(f"Extracted PPT text: {text.strip()}")
    return text.strip()

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/check', methods=['POST'])
def check():
    file = request.files.get('file')  # Get the uploaded file
    text = request.form.get('text', '').strip()  # Get the text field input

    if not text and not file:
        return "Please provide text or upload a file."

    original_text = text  # Save the original input text

    # File handling
    if file:
        filename = file.filename
        filetype = filename.split('.')[-1].lower()
        logging.debug(f"Uploaded file: {filename}, Type: {filetype}")

        if file.content_length > 10 * 1024 * 1024:  # Limit file size to 10MB
            return "The file is too large. Please upload a file smaller than 10MB."

        if filetype == 'pdf':
            text = extract_text_from_pdf(file.stream)  # Extract text from PDF
        elif filetype == 'pptx':
            text = extract_text_from_ppt(file.stream)  # Extract text from PPTX
        else:
            return "Invalid file type. Please upload a PDF or PPTX file."

        if not text.strip():
            return "The uploaded file does not contain any readable text."

    # First, perform spelling correction using LanguageTool
    text_with_correct_spelling = correct_spelling_with_tool(text)

    # Then, perform grammar correction using Grammarly's CoEdit model
    final_corrected_text = correct_grammar_with_transformers(text_with_correct_spelling)

    return render_template(
        'result.html',
        original_text=original_text,
        corrected_text=final_corrected_text
    )

@app.route('/download-pdf', methods=['POST'])
def download_pdf():
    """Generate and return a PDF with the corrected text."""
    corrected_text = request.form.get('corrected_text', '').strip()
    
    if not corrected_text:
        return "No corrected text available to save as a PDF."

    # Create an in-memory file-like object for the PDF
    pdf_buffer = io.BytesIO()
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)

    # Configure text layout
    width, height = letter
    x_margin, y_margin = 50, height - 50  # Starting point for text
    line_height = 15  # Space between lines

    # Add text to the PDF, wrapping if necessary
    pdf.setFont("Helvetica", 12)
    lines = corrected_text.split('\n')
    for line in lines:
        if y_margin - line_height < 50:  # Check if the text will overflow the page
            pdf.showPage()  # Add a new page
            y_margin = height - 50
        pdf.drawString(x_margin, y_margin, line.strip())
        y_margin -= line_height

    pdf.save()  # Finalize the PDF
    pdf_buffer.seek(0)

    # Send the generated PDF as a downloadable file
    return send_file(
        pdf_buffer,
        as_attachment=True,
        download_name="corrected_text.pdf",
        mimetype='application/pdf'
    )

if __name__ == '__main__':
    app.run(debug=True)
